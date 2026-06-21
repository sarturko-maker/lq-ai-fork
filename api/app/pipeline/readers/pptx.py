"""PPTX reader — one unit per slide, via python-pptx (MIT).

Each slide becomes one unit (``page_number`` = 1-based slide number, the
natural fit for the ``page`` field): the newline-joined text of every
shape on the slide that carries a text frame, in shape order. The
container is hardened by :func:`guard_ooxml` *before* python-pptx (lxml)
opens it.
"""

from __future__ import annotations

from io import BytesIO

from app.pipeline.readers._base import (
    OOXML_PPTX_MIME,
    ParsedDocument,
    ParserError,
    build_parsed_document,
    dist_version,
    guard_ooxml,
    ooxml_subtype,
)


class PptxReader:
    """Extracts per-slide text from a ``.pptx`` presentation."""

    parser_label = "python-pptx"
    mimes = frozenset({OOXML_PPTX_MIME})

    def sniff(self, data: bytes) -> bool:
        return ooxml_subtype(data) == "pptx"

    def read(self, data: bytes) -> ParsedDocument:
        guard_ooxml(data)
        try:
            from pptx import Presentation
        except ImportError as exc:  # pragma: no cover - install-time error
            raise ParserError("python-pptx is not installed; cannot read PPTX") from exc

        # The container passed sniff + guard_ooxml; a still-malformed deck
        # must fail closed as ParserError (-> ingestion parse_failed), not
        # escape as a bare python-pptx error (mirrors parsers._run_pymupdf).
        try:
            presentation = Presentation(BytesIO(data))
            unit_texts: list[str] = []
            for slide in presentation.slides:
                shape_texts = [
                    shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
                ]
                unit_texts.append("\n".join(shape_texts))
        except Exception as exc:
            raise ParserError(f"python-pptx failed to read PPTX: {exc}") from exc

        return build_parsed_document(
            unit_texts,
            parser=self.parser_label,
            parser_version=f"python-pptx={dist_version('python-pptx')}",
        )
